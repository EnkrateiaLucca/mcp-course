#!/usr/bin/env python3
# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""Convert remark.js presentation to PPTX format with all 86 slide states."""

import re
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# --- Configuration ---
PROJECT_ROOT = Path(__file__).parent.parent
HTML_PATH = PROJECT_ROOT / "presentation" / "presentation.html"
ASSETS_DIR = PROJECT_ROOT / "demos" / "assets-resources"
OUTPUT_PATH = PROJECT_ROOT / "presentation" / "presentation-mcp.pptx"

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.8)
CONTENT_W = Inches(11.733)

# Colors
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_BODY = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x88, 0x88, 0x88)
CODE_BG = RGBColor(0xF5, 0xF3, 0xEB)
CODE_BORDER = RGBColor(0xDD, 0xD9, 0xD2)
DEMO_GREEN = RGBColor(0x90, 0xEE, 0x90)

# Fonts
FONT_TITLE = "Droid Serif"
FONT_BODY = "Droid Serif"
FONT_CODE = "Courier New"

# Font sizes per heading level
HEADING_SIZES = {1: Pt(40), 2: Pt(32), 3: Pt(26)}
BODY_SIZE = Pt(20)
CODE_SIZE = Pt(13)
DEMO_SIZE = Pt(28)


# ============================================================
# PARSING
# ============================================================

def extract_textarea(html_path):
    """Extract the remark.js textarea content from HTML."""
    text = Path(html_path).read_text(encoding="utf-8")
    m = re.search(r'<textarea id="source">(.*?)</textarea>', text, re.DOTALL)
    if not m:
        raise ValueError("Could not find <textarea id='source'> in HTML")
    return m.group(1)


def split_slides_and_reveals(textarea_content):
    """Split textarea into 86 flattened slide states.

    1. Split by --- (slide separator)
    2. For each slide, split by -- (incremental reveal)
    3. Create cumulative states for each reveal level
    """
    # Split by --- on its own line
    raw_slides = re.split(r'\n[ \t]*---[ \t]*\n', textarea_content)

    all_states = []
    for raw in raw_slides:
        # Extract class: directive from start (may have leading newlines)
        classes = []
        content = raw
        class_match = re.search(r'^[ \t]*class:\s*(.+?)[ \t]*$', content, re.MULTILINE)
        if class_match:
            classes = [c.strip() for c in class_match.group(1).split(',')]
            content = content[:class_match.start()] + content[class_match.end():]

        # Split by -- (reveal separator), but NOT --- (already handled)
        parts = re.split(r'\n[ \t]*--[ \t]*\n', content)

        # Create cumulative states
        accumulated = ""
        for i, part in enumerate(parts):
            accumulated = accumulated + ("\n" + part if i > 0 else part)
            all_states.append({
                'classes': classes,
                'content': accumulated.strip(),
            })

    return all_states


def preprocess_content(content):
    """Normalize multi-line HTML blocks into parseable lines."""
    # Collapse multi-line demo labels into single lines
    # Pattern: <h1>\n<span style="background-color: lightgreen">\nText\n</span>\n</h1>
    # Collapse all demo labels (multi-line safe) into single-line markers
    def demo_replace(m):
        text = " ".join(m.group(1).split())  # collapse whitespace/newlines
        text = re.sub(r'<[^>]+>', '', text).strip()  # strip inner HTML tags
        return f'@@DEMO@@{text}@@/DEMO@@'

    # Match <span style="background-color: lightgreen">...</span> anywhere
    content = re.sub(
        r'<span\s+style=["\']background-color:\s*lightgreen["\']\s*>\s*(.*?)\s*</span>',
        demo_replace, content, flags=re.DOTALL
    )

    # Remove .center[ and matching ]
    content = re.sub(r'\.center\[', '', content)
    # Remove standalone ] lines that closed .center[]
    content = re.sub(r'\n\s*\]\s*\n', '\n', content)
    # Remove trailing ]
    content = re.sub(r'\n\s*\]\s*$', '', content)

    # Remove standalone HTML wrapper tags (but keep img and content)
    content = re.sub(r'^\s*</?h\d>\s*$', '', content, flags=re.MULTILINE)

    # Remove any remaining class: directives (safety net)
    content = re.sub(r'^\s*class:\s*.+$', '', content, flags=re.MULTILINE)

    return content


def parse_blocks(content):
    """Parse preprocessed content into structured blocks."""
    content = preprocess_content(content)
    blocks = []
    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # Demo label (preprocessed marker)
        demo_match = re.match(r'@@DEMO@@(.+?)@@/DEMO@@', stripped)
        if demo_match:
            text = re.sub(r'<[^>]+>', '', demo_match.group(1)).strip()
            blocks.append({'type': 'demo', 'text': text})
            i += 1
            continue

        # Code block
        if stripped.startswith('```'):
            lang = stripped[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1  # skip closing ```
            # Remove common leading whitespace
            if code_lines:
                min_indent = min((len(l) - len(l.lstrip()) for l in code_lines if l.strip()), default=0)
                code_lines = [l[min_indent:] if len(l) >= min_indent else l for l in code_lines]
            blocks.append({'type': 'code', 'lang': lang, 'text': '\n'.join(code_lines)})
            continue

        # Image tag
        img_match = re.search(r'<img\s+[^>]*src=["\']([^"\']+)["\']', stripped)
        if img_match:
            src = img_match.group(1)
            alt_m = re.search(r'alt=["\']([^"\']*)["\']', stripped)
            alt = alt_m.group(1) if alt_m else ''
            blocks.append({'type': 'image', 'src': src, 'alt': alt})
            i += 1
            continue

        # Heading (# to ######) — with or without space after #
        h_match = re.match(r'^(#{1,6})\s*(.*\S.*)', stripped)
        if h_match:
            level = len(h_match.group(1))
            text = h_match.group(2).strip()
            blocks.append({'type': 'heading', 'level': level, 'text': text})
            i += 1
            continue

        # Bullet point
        b_match = re.match(r'^[-*]\s+(.+)', stripped)
        if b_match:
            blocks.append({'type': 'bullet', 'text': b_match.group(1)})
            i += 1
            continue

        # Numbered item
        n_match = re.match(r'^\d+[.)]\s+(.+)', stripped)
        if n_match:
            blocks.append({'type': 'bullet', 'text': n_match.group(1)})
            i += 1
            continue

        # Skip empty HTML tags and wrappers
        if re.match(r'^</?(?:div|span|h\d|p|br|hr)\s*/?>$', stripped):
            i += 1
            continue

        # Plain text (strip remaining HTML tags)
        text = re.sub(r'<[^>]+>', '', stripped).strip()
        if text:
            blocks.append({'type': 'text', 'text': text})
        i += 1

    return blocks


# ============================================================
# TEXT FORMATTING HELPERS
# ============================================================

def parse_inline_markdown(text):
    """Parse markdown bold/italic/code/links into run descriptors."""
    runs = []
    pos = 0

    # Combined pattern for bold, italic, inline code, links
    pattern = re.compile(
        r'\*\*(.+?)\*\*'       # **bold**
        r'|\*([^*]+?)\*'       # *italic*
        r'|`([^`]+?)`'         # `code`
        r'|\[([^\]]+?)\]\([^\)]+?\)'  # [link text](url)
    )

    for m in pattern.finditer(text):
        # Add plain text before this match
        if m.start() > pos:
            runs.append({'text': text[pos:m.start()], 'bold': False, 'italic': False, 'code': False})

        if m.group(1) is not None:  # bold
            runs.append({'text': m.group(1), 'bold': True, 'italic': False, 'code': False})
        elif m.group(2) is not None:  # italic
            runs.append({'text': m.group(2), 'bold': False, 'italic': True, 'code': False})
        elif m.group(3) is not None:  # code
            runs.append({'text': m.group(3), 'bold': False, 'italic': False, 'code': True})
        elif m.group(4) is not None:  # link
            runs.append({'text': m.group(4), 'bold': False, 'italic': False, 'code': False})

        pos = m.end()

    # Remaining text
    if pos < len(text):
        runs.append({'text': text[pos:], 'bold': False, 'italic': False, 'code': False})

    if not runs:
        runs.append({'text': text, 'bold': False, 'italic': False, 'code': False})

    return runs


def add_formatted_text(paragraph, text, base_font=FONT_BODY, base_size=BODY_SIZE,
                       base_color=TEXT_BODY, base_bold=False):
    """Add markdown-formatted text runs to a paragraph."""
    inline_runs = parse_inline_markdown(text)
    for rd in inline_runs:
        run = paragraph.add_run()
        run.text = rd['text']
        run.font.name = FONT_CODE if rd['code'] else base_font
        run.font.size = Pt(base_size.pt - 2) if rd['code'] else base_size
        run.font.color.rgb = base_color
        run.font.bold = rd['bold'] or base_bold
        run.font.italic = rd['italic']


# ============================================================
# IMAGE HELPERS
# ============================================================

def resolve_image(src):
    """Resolve image src (relative to presentation/) to absolute path."""
    if src.startswith('../'):
        path = PROJECT_ROOT / src[3:]
    elif src.startswith('./'):
        path = PROJECT_ROOT / "presentation" / src[2:]
    else:
        path = PROJECT_ROOT / "presentation" / src
    return path.resolve()


def ensure_compatible_image(image_path):
    """Convert WEBP or other unsupported formats to PNG. Returns path to use."""
    try:
        with Image.open(image_path) as img:
            if img.format in ('BMP', 'GIF', 'JPEG', 'PNG', 'TIFF'):
                return image_path
            # Convert to PNG
            png_path = image_path.with_suffix('.converted.png')
            img.save(png_path, 'PNG')
            return png_path
    except Exception:
        return image_path


def fit_image(image_path, max_w_inches, max_h_inches):
    """Calculate image dimensions maintaining aspect ratio. Returns (w_emu, h_emu)."""
    try:
        with Image.open(image_path) as img:
            iw, ih = img.size
    except Exception:
        return Inches(max_w_inches * 0.8), Inches(max_h_inches * 0.5)

    aspect = iw / ih
    max_w = Inches(max_w_inches)
    max_h = Inches(max_h_inches)

    # Scale to fit
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)

    return w, h


# ============================================================
# SLIDE CREATION
# ============================================================

def make_slide(prs, state):
    """Create one PPTX slide from a parsed state."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_WHITE

    classes = state['classes']
    blocks = parse_blocks(state['content'])

    centered = 'center' in classes
    middle = 'middle' in classes

    if not blocks:
        return slide

    # Estimate total content height for vertical centering
    total_h = estimate_height(blocks)

    if middle:
        start_y = max(0.4, (7.5 - total_h) / 2)
    else:
        start_y = 0.5

    y = start_y

    for block in blocks:
        y = render_block(slide, block, y, centered)

    return slide


def estimate_height(blocks):
    """Rough estimate of total content height in inches."""
    h = 0
    for b in blocks:
        t = b['type']
        if t == 'heading':
            h += {1: 0.9, 2: 0.8, 3: 0.7}.get(b['level'], 0.6)
        elif t == 'image':
            h += 4.5
        elif t == 'demo':
            h += 1.5
        elif t == 'code':
            lines = b['text'].count('\n') + 1
            h += max(1.0, lines * 0.24 + 0.4)
        elif t in ('bullet', 'text'):
            h += 0.5
    return h


def render_block(slide, block, y, centered):
    """Render a content block onto the slide at vertical position y. Returns new y."""
    btype = block['type']

    if btype == 'heading':
        return render_heading(slide, block, y, centered)
    elif btype == 'image':
        return render_image(slide, block, y)
    elif btype == 'demo':
        return render_demo(slide, block, y)
    elif btype == 'code':
        return render_code(slide, block, y)
    elif btype == 'bullet':
        return render_bullet(slide, block, y, centered)
    elif btype == 'text':
        return render_text(slide, block, y, centered)

    return y


def render_heading(slide, block, y, centered):
    level = block['level']
    size = HEADING_SIZES.get(level, BODY_SIZE)
    h_inches = {1: 0.9, 2: 0.8, 3: 0.7}.get(level, 0.6)

    txBox = slide.shapes.add_textbox(MARGIN_X, Inches(y), CONTENT_W, Inches(h_inches))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

    add_formatted_text(p, block['text'], base_font=FONT_TITLE, base_size=size,
                       base_color=TEXT_DARK, base_bold=(level <= 2))

    return y + h_inches


def render_image(slide, block, y):
    image_path = resolve_image(block['src'])

    if not image_path.exists():
        # Placeholder text
        txBox = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(9), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[Image: {block.get('alt', block['src'])}]"
        run.font.name = FONT_BODY
        run.font.size = Pt(16)
        run.font.color.rgb = TEXT_GRAY
        return y + 0.6

    # Convert unsupported formats (e.g. WEBP) to PNG
    image_path = ensure_compatible_image(image_path)

    remaining_h = max(2.0, 7.0 - y)
    img_w, img_h = fit_image(image_path, 11.5, remaining_h)

    # Center horizontally
    x = int((SLIDE_WIDTH - img_w) / 2)
    slide.shapes.add_picture(str(image_path), x, Inches(y), img_w, img_h)

    actual_h = img_h / 914400  # EMU to inches
    return y + actual_h + 0.2


def render_demo(slide, block, y):
    """Green highlighted demo/whiteboard label."""
    text = block['text']

    box_w = Inches(10)
    box_h = Inches(1.4)
    x = int((SLIDE_WIDTH - box_w) / 2)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y), box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DEMO_GREEN
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    # Vertical center
    tf._txBody.bodyPr.set('anchor', 'ctr')

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TITLE
    run.font.size = DEMO_SIZE
    run.font.color.rgb = TEXT_DARK
    run.font.bold = True

    return y + 1.7


def render_code(slide, block, y):
    code = block['text']
    lines = code.count('\n') + 1
    code_h = max(1.0, lines * 0.24 + 0.3)

    # Clamp to available space
    remaining = 7.0 - y
    if code_h > remaining:
        code_h = max(1.0, remaining - 0.1)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(y), CONTENT_W, Inches(code_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CODE_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = code
    run.font.name = FONT_CODE
    run.font.size = CODE_SIZE
    run.font.color.rgb = TEXT_BODY

    return y + code_h + 0.15


def render_bullet(slide, block, y, centered):
    txBox = slide.shapes.add_textbox(
        Inches(1.2) if not centered else MARGIN_X,
        Inches(y),
        Inches(10.9) if not centered else CONTENT_W,
        Inches(0.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

    # Add bullet prefix then formatted text
    run = p.add_run()
    run.text = "•  "
    run.font.name = FONT_BODY
    run.font.size = BODY_SIZE
    run.font.color.rgb = TEXT_BODY

    add_formatted_text(p, block['text'])

    return y + 0.48


def render_text(slide, block, y, centered):
    txBox = slide.shapes.add_textbox(MARGIN_X, Inches(y), CONTENT_W, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

    add_formatted_text(p, block['text'])

    return y + 0.48


# ============================================================
# MAIN
# ============================================================

def main():
    print(f"Reading {HTML_PATH}")
    textarea = extract_textarea(HTML_PATH)

    states = split_slides_and_reveals(textarea)
    print(f"Parsed {len(states)} slide states (expected 86)")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, state in enumerate(states):
        make_slide(prs, state)
        if (i + 1) % 10 == 0:
            print(f"  Created slide {i + 1}/{len(states)}")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved {len(states)} slides to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
